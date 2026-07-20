from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("marisco_declarative_intake.pptx")
LOGO = Path(__file__).parent / "nbs" / "assets" / "logo.png"

NAVY = "071A2E"
NAVY_2 = "0C2740"
TEAL = "19D3C5"
MINT = "C8F28A"
SKY = "72B9FF"
CORAL = "FF7D68"
WHITE = "F5F8FC"
MUTED = "A9B9C8"
LINE = "23425A"
INK = "0A1724"


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill, radius=False, line=None, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
         margin=0.05, spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    para = frame.paragraphs[0]
    para.alignment = align
    para.space_after = Pt(0)
    para.line_spacing = spacing
    run = para.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def label(slide, value, x, y, w, color=TEAL):
    return text(slide, value.upper(), x, y, w, 0.23, size=10, color=color, bold=True,
                font="Aptos Display", margin=0)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def footer(slide, page):
    line(slide, 0.55, 7.05, 12.78, 7.05, LINE, 0.8)
    text(slide, "MARIS  |  Declarative Intake", 0.62, 7.13, 3.1, 0.18,
         size=8.5, color=MUTED, bold=True, margin=0)
    text(slide, f"{page:02d}", 12.25, 7.11, 0.45, 0.2,
         size=9, color=TEAL, bold=True, align=PP_ALIGN.RIGHT, margin=0)


def title(slide, eyebrow, headline, subhead, page):
    label(slide, eyebrow, 0.62, 0.42, 4.8)
    text(slide, headline, 0.62, 0.72, 11.8, 0.58, size=31, color=WHITE,
         bold=True, font="Aptos Display", margin=0)
    text(slide, subhead, 0.64, 1.38, 11.5, 0.34, size=14.5, color=MUTED, margin=0)
    footer(slide, page)


def chip(slide, value, x, y, w, fill=NAVY_2, accent=TEAL):
    rect(slide, x, y, w, 0.34, fill, radius=True, line=accent)
    text(slide, value, x + 0.12, y + 0.075, w - 0.24, 0.16, size=9.5,
         color=WHITE, bold=True, margin=0, align=PP_ALIGN.CENTER)


def code_block(slide, lines, x, y, w, h, accent=TEAL):
    rect(slide, x, y, w, h, "06111E", radius=True, line=LINE)
    rect(slide, x, y, 0.07, h, accent, radius=False, line=accent)
    text(slide, lines, x + 0.2, y + 0.16, w - 0.35, h - 0.25, size=12,
         color="D9E7F2", font="Consolas", margin=0, spacing=0.95)


def add_logo(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.83), Inches(0.34), width=Inches(0.82))


def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    rect(slide, 9.85, 0.0, 3.48, 7.5, "123B58", radius=True, line="123B58", transparency=15)
    add_logo(slide)
    label(slide, "MARIS v1.5 / engineering review", 0.65, 0.53, 4.6)
    text(slide, "Provider-specific mess\nstops at the boundary.", 0.62, 0.92, 8.2, 1.42,
         size=38, color=WHITE, bold=True, font="Aptos Display", margin=0)
    text(slide, "Declarative Intake makes every dataset extension a local boundary + contract change — not a new core-pipeline branch.",
         0.66, 2.47, 7.45, 0.56, size=17, color=MUTED, margin=0)
    chip(slide, "NO HANDLER MIGRATION", 0.65, 3.22, 1.88)
    chip(slide, "3 DATASETS VERIFIED", 2.67, 3.22, 1.82, fill=NAVY_2, accent=MINT)

    # Arrows behind the boundary nodes.
    line(slide, 3.98, 5.02, 4.32, 5.02, TEAL, 2.0)
    line(slide, 7.22, 5.02, 7.56, 5.02, TEAL, 2.0)
    for x in (4.25, 7.49):
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(4.93), Inches(0.18), Inches(0.18))
        tri.rotation = 90
        tri.fill.solid(); tri.fill.fore_color.rgb = rgb(TEAL); tri.line.color.rgb = rgb(TEAL)

    cards = [
        (0.65, "01", "BOUNDARY LOADER", "ZIP / XLSX\nbroken headers\ndate parsing", CORAL),
        (4.36, "02", "YAML CONTRACT", "columns · melt · LUTs\ndefaults · metadata", MINT),
        (7.60, "03", "SHARED ENGINE", "provider-blind gates\nNetCDF output", SKY),
    ]
    widths = [3.18, 2.72, 3.25]
    for (x, num, heading, body, accent), width in zip(cards, widths):
        rect(slide, x, 4.06, width, 1.64, NAVY_2, radius=True, line=LINE)
        text(slide, num, x + 0.18, 4.25, 0.35, 0.28, size=13, color=accent, bold=True, margin=0)
        text(slide, heading, x + 0.18, 4.58, width - 0.32, 0.24, size=13, color=WHITE, bold=True, margin=0)
        text(slide, body, x + 0.18, 4.96, width - 0.3, 0.58, size=12.5, color=MUTED, margin=0)
    text(slide, "The review question: do provider quirks stay in the Loader, while meaning stays declarative?",
         0.66, 6.13, 10.9, 0.32, size=14, color=TEAL, bold=True, margin=0)
    footer(slide, 1)


def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    title(slide, "Concrete practical win", "JOIS avoids an SQL import crash in two YAML lines.",
          "The provider file has no detection field. The contract supplies the required semantic default at the right point in the flow.", 2)

    # raw provider image
    rect(slide, 0.63, 2.03, 3.62, 3.92, NAVY_2, radius=True, line=LINE)
    label(slide, "after Boundary Loader", 0.88, 2.27, 2.6, color=CORAL)
    text(slide, "Provider-shaped data", 0.88, 2.58, 2.9, 0.28, size=18, color=WHITE, bold=True, margin=0)
    code_block(slide,
        "Cruise       JOIS 2021\nStation      CB4\nLatitude     75.000833\nLongitude   -150.001333\nI129_at_kg  6.40e+08\n\nRAW_DETECTION_PRESENT=False",
        0.88, 3.08, 3.1, 1.92, accent=CORAL)
    text(slide, "Physical file cleanup is complete.\nNo MARIS semantic field is invented here.",
         0.9, 5.24, 3.04, 0.42, size=11.5, color=MUTED, margin=0)

    # central contract
    rect(slide, 4.62, 2.03, 3.74, 3.92, "102E47", radius=True, line="24506C")
    label(slide, "YAML contract", 4.9, 2.27, 2.3, color=MINT)
    text(slide, "Declare the SQL semantic", 4.9, 2.58, 3.0, 0.28, size=18, color=WHITE, bold=True, margin=0)
    code_block(slide, "measurement_defaults:\n  detection: \"=\"", 4.9, 3.12, 3.12, 0.9, accent=MINT)
    text(slide, "Applied post-melt\nApplied before Gate 2\nExisting non-null values are preserved", 4.92, 4.38, 3.1, 0.8,
         size=13, color=MUTED, margin=0)
    chip(slide, "NO CUSTOM CALLBACK", 4.92, 5.26, 2.72, fill="16354D", accent=MINT)

    # final result
    rect(slide, 8.73, 2.03, 3.95, 3.92, NAVY_2, radius=True, line=LINE)
    label(slide, "after declarative preflight", 9.0, 2.27, 3.0, color=TEAL)
    text(slide, "76 valid SQL-ready records", 9.0, 2.58, 3.25, 0.28, size=18, color=WHITE, bold=True, margin=0)
    code_block(slide, "detection\n        =\n        =\n        =\n        =\n\nDETECTION_COUNTS={'=': 76}", 9.0, 3.08, 3.3, 1.92, accent=TEAL)
    text(slide, "CLI: Gate 1 + Gate 2 passed\nSEAWATER: 76 rows × 17 cols", 9.02, 5.24, 3.2, 0.42,
         size=11.5, color=MUTED, margin=0)
    footer(slide, 2)


def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    title(slide, "Safety before lossy output", "Gate 2 exposes defects before sanitization can hide them.",
          "Validation is deliberately positioned before coordinate/time cleanup and NetCDF encoding.", 3)

    # gate path
    rect(slide, 0.63, 2.08, 4.0, 4.45, NAVY_2, radius=True, line=LINE)
    label(slide, "validation path", 0.92, 2.33, 2.4, color=SKY)
    steps = [
        ("GATE 1", "Pydantic validates YAML shape", SKY),
        ("TRANSFORM", "Map, parse, melt, remap", TEAL),
        ("GATE 2", "Check required columns, LAT/LON/TIME", CORAL),
        ("OUTPUT", "Sanitize + encode only after pass", MINT),
    ]
    y = 2.74
    for idx, (tag, desc, accent) in enumerate(steps):
        rect(slide, 0.95, y, 3.34, 0.62, "102E47", radius=True, line="214861")
        text(slide, tag, 1.13, y + 0.13, 0.88, 0.18, size=10.5, color=accent, bold=True, margin=0)
        text(slide, desc, 2.0, y + 0.11, 2.02, 0.22, size=11.5, color=WHITE, margin=0)
        if idx < len(steps) - 1:
            line(slide, 2.62, y + 0.62, 2.62, y + 0.79, TEAL, 1.4)
        y += 0.88
    text(slide, "Key principle: do not let downstream cleanup\ndiscard evidence of upstream defects.",
         0.96, 6.04, 3.2, 0.34, size=11.5, color=MUTED, margin=0)

    # actual terminal output
    rect(slide, 4.98, 2.08, 7.7, 4.45, "06111E", radius=True, line=LINE)
    label(slide, "raw probe / actual Gate 2 diagnostic", 5.26, 2.33, 3.6, color=CORAL)
    text(slide, "Fram Strait 2024 Raw Probe", 5.26, 2.62, 4.1, 0.24, size=16, color=WHITE, bold=True, margin=0)
    code_block(slide,
        "Group 'SEAWATER':\n- LAT contains 4 null value(s).\n- LON contains 4 null value(s).\n- TIME contains 376 null value(s).\n\nEncoding would be unsafe here because downstream\ntime/coordinate rails can discard rows and mask\nthe upstream defect.",
        5.25, 3.04, 4.55, 2.33, accent=CORAL)
    rect(slide, 10.15, 3.04, 2.15, 2.33, "102E47", radius=True, line="214861")
    text(slide, "ACTION\n\nRepair the\nBoundary Loader\n\n—not the\nshared pipeline.", 10.4, 3.34, 1.7, 1.67,
         size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    text(slide, "The diagnostic prints a ready-to-adapt custom Loader skeleton and exits non-zero.",
         5.27, 5.71, 6.66, 0.28, size=12.5, color=TEAL, bold=True, margin=0)
    footer(slide, 3)


def slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    title(slide, "Review with confidence", "Proven on three formats. Governed for the fourth.",
          "The review focus moves from bespoke callback code to explicit boundaries, contracts, and evidence.", 4)

    datasets = [
        ("CERN", "524 × 17", "mixed dates + radionuclides", SKY),
        ("JOIS", "76 × 17", "ZIP/XLSX + SQL default", TEAL),
        ("FRAM STRAIT", "377 × 17", "merged headers + scales", MINT),
    ]
    x = 0.64
    for name, metric, detail, accent in datasets:
        rect(slide, x, 2.12, 2.5, 1.52, NAVY_2, radius=True, line=LINE)
        text(slide, name, x + 0.22, 2.35, 2.0, 0.2, size=11, color=accent, bold=True, margin=0)
        text(slide, metric, x + 0.22, 2.67, 2.12, 0.35, size=23, color=WHITE, bold=True, margin=0)
        text(slide, detail, x + 0.22, 3.14, 2.08, 0.2, size=10.5, color=MUTED, margin=0)
        x += 2.68

    rect(slide, 8.86, 2.12, 3.82, 1.52, "102E47", radius=True, line="214861")
    label(slide, "runtime compatibility", 9.14, 2.35, 2.4, color=TEAL)
    text(slide, "Legacy callback logic\nremains untouched", 9.14, 2.68, 3.12, 0.5, size=17, color=WHITE, bold=True, margin=0)
    text(slide, "Only an nbdev-generated Docs URL changed.", 9.14, 3.23, 3.1, 0.16, size=9.5, color=MUTED, margin=0)

    rect(slide, 0.64, 4.08, 12.04, 2.18, "102E47", radius=True, line="214861")
    label(slide, "INTAKE_GOVERNANCE.md / review contract", 0.94, 4.34, 4.0, color=MINT)
    governance = [
        ("1", "Keep file physics in the Loader", "ZIPs, sheets, broken cells, scale cleanup"),
        ("2", "Keep meaning in YAML", "mappings, melt rules, LUTs, measurement defaults"),
        ("3", "Keep the core provider-blind", "mandatory gates, no dataset branches"),
    ]
    x = 0.94
    for number, heading, body in governance:
        text(slide, number, x, 4.84, 0.32, 0.25, size=14, color=MINT, bold=True, margin=0)
        text(slide, heading, x + 0.38, 4.78, 3.15, 0.24, size=14, color=WHITE, bold=True, margin=0)
        text(slide, body, x + 0.38, 5.2, 3.15, 0.38, size=10.5, color=MUTED, margin=0)
        x += 3.92

    rect(slide, 0.64, 6.49, 12.04, 0.35, "16354D", radius=True, line=TEAL)
    text(slide, "REVIEW DECISION  |  Approve the boundary discipline: future datasets add contracts and loaders, not core branches.",
         0.87, 6.58, 11.55, 0.15, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    footer(slide, 4)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    prs.save(OUT)
    print(f"Created {OUT}")


if __name__ == "__main__":
    build_deck()
